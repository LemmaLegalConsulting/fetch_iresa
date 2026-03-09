#!/usr/bin/env python3
"""Simple converter: Markdown slide -> PPTX
Reads one markdown slide file and produces a multi-slide PowerPoint.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import re
from pathlib import Path

MD_PATH = Path("../slides/followup_question_judging_slide.md")
OUT_PATH = Path("../slides/followup_question_judging.pptx")


def parse_markdown(md_text: str):
    # Split into sections by a line with three dashes
    parts = re.split(r"\n-{3,}\n", md_text)
    sections = []
    for p in parts:
        p = p.strip()
        if not p:
            continue
        sections.append(p)
    return sections


def extract_title_and_subtitle(sec: str):
    lines = [l.rstrip() for l in sec.splitlines() if l.strip()]
    title = ""
    subtitle = ""
    if lines:
        if lines[0].startswith("#"):
            title = lines[0].lstrip("# ")
            subtitle = "\n".join(lines[1:3]).strip()
        else:
            title = lines[0]
            subtitle = "\n".join(lines[1:3]).strip()
    return title, subtitle


def extract_sections(sec: str):
    # Find top-level headings (##) and collect following bullets/paragraphs
    lines = [l.rstrip() for l in sec.splitlines()]
    cur_h = None
    cur_b = []
    out = []
    for ln in lines:
        if ln.startswith('##'):
            if cur_h:
                out.append((cur_h, '\n'.join(cur_b).strip()))
            cur_h = ln.lstrip('# ').strip()
            cur_b = []
        else:
            cur_b.append(ln)
    if cur_h:
        out.append((cur_h, '\n'.join(cur_b).strip()))
    return out


def bullets_from_block(block: str):
    bullets = []
    for ln in block.splitlines():
        ln = ln.strip()
        if not ln:
            continue
        if ln.startswith('- '):
            bullets.append(ln[2:].strip())
        else:
            bullets.append(ln)
    return bullets


if __name__ == '__main__':
    md_file = (Path(__file__).parent / MD_PATH).resolve()
    out_file = (Path(__file__).parent / OUT_PATH).resolve()

    md_text = md_file.read_text(encoding='utf8')
    parts = parse_markdown(md_text)

    prs = Presentation()
    # Title slide: first part
    title, subtitle = extract_title_and_subtitle(parts[0])
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        tx = slide.placeholders[1].text_frame
        tx.clear()
        p = tx.paragraphs[0]
        p.font.size = Pt(18)
        p.text = subtitle

    # Remaining parts: parse sections and create slides
    for sec in parts[1:]:
        secs = extract_sections(sec)
        for heading, body in secs:
            # Create a slide for each heading
            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = heading
            tf = s.placeholders[1].text_frame
            tf.clear()
            bullets = bullets_from_block(body)
            for i, b in enumerate(bullets):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = b
                p.level = 0
                p.font.size = Pt(16)

    # Speaker notes (try to find Speaker notes: section)
    speaker = None
    m = re.search(r"\*\*Speaker notes:\*\*(.*)$", md_text, flags=re.S | re.I)
    if m:
        speaker = m.group(1).strip()
        # attach to last slide
        if prs.slides:
            notes_slide = prs.slides[-1].notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = speaker

    prs.save(out_file)
    print(f"Saved {out_file}")
